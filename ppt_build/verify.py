"""Sanity-check the generated PPTX: per-slide text dump + shape bounds +
text-overflow heuristic."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu

p = Presentation("EmoMoconvq_spotlight.pptx")

SLIDE_W = p.slide_width
SLIDE_H = p.slide_height
print(f"Slide size: {Emu(SLIDE_W).inches:.2f} x {Emu(SLIDE_H).inches:.2f} in")
print(f"# slides:   {len(p.slides)}\n")

for i, slide in enumerate(p.slides, 1):
    print(f"================ Slide {i} ================")
    for j, sh in enumerate(slide.shapes):
        try:
            l, t = sh.left, sh.top
            w, h = sh.width, sh.height
        except Exception:
            continue
        bounds = (Emu(l).inches, Emu(t).inches,
                  Emu(w).inches, Emu(h).inches)
        kind = sh.shape_type
        txt = ""
        if sh.has_text_frame:
            tf = sh.text_frame
            txt = " | ".join(p.text for p in tf.paragraphs if p.text).strip()
            if len(txt) > 110:
                txt = txt[:110] + "..."
        flag = ""
        if (bounds[0] + bounds[2] > Emu(SLIDE_W).inches + 0.01 or
            bounds[1] + bounds[3] > Emu(SLIDE_H).inches + 0.01 or
            bounds[0] < -0.01 or bounds[1] < -0.01):
            flag = "  OUT-OF-BOUNDS"
        if txt:
            print(f"  [{j:02d}] {kind!s:>14} "
                  f"({bounds[0]:.2f},{bounds[1]:.2f}) "
                  f"{bounds[2]:.2f}x{bounds[3]:.2f}{flag}")
            print(f"        \"{txt}\"")
        else:
            print(f"  [{j:02d}] {kind!s:>14} "
                  f"({bounds[0]:.2f},{bounds[1]:.2f}) "
                  f"{bounds[2]:.2f}x{bounds[3]:.2f}{flag}")
